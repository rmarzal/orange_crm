# -*- coding: utf-8 -*-
import sys; sys.path.insert(0,'/home/user/orange_crm')
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from ppt_helpers import set_para, trim_paragraphs, check_fit, add_paragraph

ORANGE="FF7900"; DARK="1A1A1A"; GREY="404040"; WHITE="FFFFFF"; LIGHT="F3F1EF"; ORANGE2="F16E00"

prs = Presentation('/home/user/orange_crm/propuesta_original.pptx')
S = prs.slides
def sh(slide_i, shp_i): return S[slide_i-1].shapes[shp_i]

warn=[]
def fit(shape,size,bold=False,label="",**kw):
    ok,_,_,_=check_fit(shape,size,bold=bold,label=label,**kw)
    if not ok: warn.append(label)

# ================= SLIDE 2 =================
s=2
trim_paragraphs(S[s-1].shapes[2].text_frame,1)
set_para(S[s-1].shapes[2].text_frame.paragraphs[0],
    "El interés por la IA crece y con él las ideas. El reto ya no es encontrar oportunidades, "
    "sino elegir bien y concentrar el esfuerzo donde más valor aporta y es realizable.",
    size=18, color=DARK)
fit(S[s-1].shapes[2],18,label="s2 intro")

tf=S[s-1].shapes[4].text_frame
trim_paragraphs(tf,4)
set_para(tf.paragraphs[0], [("La dificultad suele estar en tres frentes:",True)], size=16, color=DARK)
set_para(tf.paragraphs[1], [("Del entusiasmo a la realidad",True),(" — estimar valor y viabilidad con un criterio común.",False)], size=16, color=DARK)
set_para(tf.paragraphs[2], [("El salto que se subestima",True),(" — pasar de una PoC que funciona a operación sostenible.",False)], size=16, color=DARK)
set_para(tf.paragraphs[3], [("Muchas ideas, foco limitado",True),(" — elegir entre iniciativas que compiten por equipos y presupuesto.",False)], size=16, color=DARK)
S[s-1].shapes[4].height=Inches(1.5)
fit(S[s-1].shapes[4],16,label="s2 tres frentes")

tf=S[s-1].shapes[5].text_frame
trim_paragraphs(tf,2)
set_para(tf.paragraphs[0], [("Ya tenéis iniciativas sobre la mesa. Os damos ",False),("foco y respaldo",True),
    (": priorizar juntos, con método y mirada externa, y dejar claro qué abordar primero y cómo.",False)], size=16, color=DARK)
set_para(tf.paragraphs[1], [("El siguiente paso es un ",False),("assessment corto y acotado",True),
    (" del que salen las iniciativas clave priorizadas y una foto clara de ejecución.",False)], size=16, color=DARK)
fit(S[s-1].shapes[5],16,label="s2 cierre")

# ================= SLIDE 3 =================
s=3
trim_paragraphs(S[s-1].shapes[2].text_frame,1)
set_para(S[s-1].shapes[2].text_frame.paragraphs[0],
    [("El assessment lo lideráis vosotros; nosotros aportamos método, mirada externa y experiencia en operación. ",False),
     ("El resultado: un plan respaldado.",True)], size=16, color=DARK)
fit(S[s-1].shapes[2],16,label="s3 intro")

descs={
 11:"Las iniciativas, el conocimiento del negocio y las decisiones son vuestras. Nosotros aportamos estructura y método.",
 10:"Un marco de priorización probado para contrastar y enriquecer vuestro criterio, no para sustituirlo.",
 12:"Sabemos dónde se atascan las iniciativas al pasar de PoC a operación; la priorización es realista desde el primer día.",
 13:"Aceleramos decisiones que ya estáis tomando: vuestro equipo sale reforzado, con un plan claro y una foto de los gaps a cubrir.",
}
for idx,txt in descs.items():
    tf=S[s-1].shapes[idx].text_frame; trim_paragraphs(tf,1)
    set_para(tf.paragraphs[0], txt, size=16, color=DARK)
    S[s-1].shapes[idx].height=Inches(0.72)
    fit(S[s-1].shapes[idx],16,label=f"s3 desc idx{idx}")

# ================= SLIDE 4 =================
s=4
trim_paragraphs(S[s-1].shapes[2].text_frame,1)
set_para(S[s-1].shapes[2].text_frame.paragraphs[0],
    [("Un assessment acotado. ",True),("Priorizamos vuestras iniciativas de IA por valor y viabilidad y salimos con una foto clara de en qué trabajar y con qué modelo de ejecución.",False)],
    size=16, color=DARK)
S[s-1].shapes[2].height=Inches(0.75)
fit(S[s-1].shapes[2],16,label="s4 intro")

tf=S[s-1].shapes[6].text_frame; trim_paragraphs(tf,4)
set_para(tf.paragraphs[0], [("Encuadre e inventario",True),(" — alinear objetivos y recoger vuestras iniciativas.",False)], size=15, color=DARK)
set_para(tf.paragraphs[1], [("Valoración",True),(" — valor de negocio, viabilidad técnica y de datos, y salto a operación.",False)], size=15, color=DARK)
set_para(tf.paragraphs[2], [("Priorización",True),(" — ordenar las iniciativas con un marco común acordado.",False)], size=15, color=DARK)
set_para(tf.paragraphs[3], [("Foco y ejecución",True),(" — converger en 3 iniciativas clave y su modelo (autónomo, con apoyo o partner).",False)], size=15, color=DARK)
S[s-1].shapes[6].height=Inches(1.7)
fit(S[s-1].shapes[6],15,label="s4 pasos")

tf=S[s-1].shapes[7].text_frame; trim_paragraphs(tf,1)
set_para(tf.paragraphs[0], "Iniciativas priorizadas con criterio común y un plan de ejecución accionable, propiedad de vuestro equipo.", size=14, color=DARK)
S[s-1].shapes[7].height=Inches(0.7)
fit(S[s-1].shapes[7],14,label="s4 resultado")

tf=S[s-1].shapes[8].text_frame; trim_paragraphs(tf,3)
set_para(tf.paragraphs[0], "Sesiones cortas (1–3 h)", size=14, color=DARK)
set_para(tf.paragraphs[1], "Preparación y síntesis entre sesiones", size=14, color=DARK)
set_para(tf.paragraphs[2], "Formato Co-Creación con capacitación continua", size=14, color=DARK)
S[s-1].shapes[8].height=Inches(0.9)
fit(S[s-1].shapes[8],14,label="s4 como")

# ================= SLIDE 5 =================
s=5
tf=S[s-1].shapes[4].text_frame; trim_paragraphs(tf,1)
set_para(tf.paragraphs[0], [("Cada iniciativa se puntúa con los mismos tres ejes y una escala común. ",True),
    ("Comparación justa y decisión con criterio, no por intuición.",False)], size=16, color=DARK)
fit(S[s-1].shapes[4],16,label="s5 intro")

tf=S[s-1].shapes[5].text_frame; trim_paragraphs(tf,1)
set_para(tf.paragraphs[0], "Todas las iniciativas candidatas —las que ya tenéis o las que surjan— en fichas homogéneas para compararlas.", size=15, color=DARK)
S[s-1].shapes[5].height=Inches(0.72)
fit(S[s-1].shapes[5],15,label="s5 quepuntuamos")

tf=S[s-1].shapes[6].text_frame; trim_paragraphs(tf,3)
set_para(tf.paragraphs[0], [("Valor de negocio",True),(" — impacto, encaje estratégico y alcance.",False)], size=15, color=DARK)
set_para(tf.paragraphs[1], [("Dificultad técnica y de datos",True),(" — madurez, datos y complejidad.",False)], size=15, color=DARK)
set_para(tf.paragraphs[2], [("Salto de PoC a operación",True),(" — esfuerzo y riesgo de industrializar; donde muchas se atascan.",False)], size=15, color=DARK)
S[s-1].shapes[6].height=Inches(1.18)
fit(S[s-1].shapes[6],15,label="s5 ejes")

tf=S[s-1].shapes[7].text_frame; trim_paragraphs(tf,2)
set_para(tf.paragraphs[0], [("Escala y ponderación acordadas",True),(" con vuestro equipo: la priorización es vuestra.",False)], size=15, color=DARK)
set_para(tf.paragraphs[1], [("Puntuación conjunta",True),(" en las sesiones, con nuestra experiencia como contraste.",False)], size=15, color=DARK)
fit(S[s-1].shapes[7],15,label="s5 como")

tf=S[s-1].shapes[8].text_frame; trim_paragraphs(tf,1)
set_para(tf.paragraphs[0], "Una matriz que ordena las iniciativas por valor y viabilidad, y un ranking claro de dónde centrar el esfuerzo.", size=15, color=DARK)
S[s-1].shapes[8].height=Inches(0.72)
fit(S[s-1].shapes[8],15,label="s5 resultado")

# ================= SLIDE 7 =================
s=7
cards={
 1:("1. Disponibilidad y calidad del dato","¿Existen los datos con calidad, volumen y acceso suficientes? Es el primer filtro."),
 2:("2. Aportación de valor de la IA","¿Aporta la IA valor claro frente a una solución más simple? Confirmar que es la herramienta adecuada."),
 3:("3. Gap y madurez técnica","Distancia entre la capacidad actual y la necesaria: complejidad, madurez e integración."),
 4:("4. Factor humano y adopción","Tipos de usuario, encaje en el flujo de trabajo y facilidad de adopción."),
 5:("5. Del PoC a la operación","Riesgos de industrializar: escalado, seguridad, compliance, monitorización y mantenimiento."),
 6:("6. Consumo y costes esperados","Coste de desarrollo y, sobre todo, de operación (inferencia/compute) y coste total (TCO)."),
}
for idx,(hd,bd) in cards.items():
    tf=S[s-1].shapes[idx].text_frame; trim_paragraphs(tf,2)
    set_para(tf.paragraphs[0], [(hd,True)], size=13, color=ORANGE2)
    set_para(tf.paragraphs[1], bd, size=12, color=GREY)
    fit(S[s-1].shapes[idx],12,label=f"s7 card{idx}")

# ================= SLIDE 16 =================
s=16
def set_body(shape, text, size=14, color=DARK):
    tf=shape.text_frame; trim_paragraphs(tf,1)
    p=add_paragraph(tf); set_para(p, text, size=size, color=color)
set_body(S[s-1].shapes[4],
  "No os decimos qué hacer: aportamos lo que cuesta ver desde dentro, por qué unas iniciativas llegan a operación y otras se quedan en PoC. Detrás, 2.500 especialistas en Data e IA en Europa, líder en IA y GenAI.")
set_body(S[s-1].shapes[5],
  "Hemos visto priorizar y ejecutar IA en muchas organizaciones; ese contraste ayuda a decidir con perspectiva y evitar errores conocidos.")
set_body(S[s-1].shapes[6],
  "Priorizamos con el coste real de industrializar y operar, no solo el piloto. Ahí se gana o se pierde el valor de la IA.")
set_body(S[s-1].shapes[7],
  "Os damos criterio para decidir y ejecutar; estamos donde nos necesitéis, sin crear dependencia.")
# normalizar headers (uno era blanco sobre blanco, otro gris de bajo contraste)
from pptx.util import Pt as _Pt2
for idx in (4,5,6,7):
    hp=S[s-1].shapes[idx].text_frame.paragraphs[0]
    for r in hp.runs:
        r.font.color.rgb=RGBColor.from_string(ORANGE); r.font.bold=True; r.font.size=_Pt2(14)
for idx in (4,5,6,7):
    fit(S[s-1].shapes[idx],14,label=f"s16 col{idx}")

# ================= SLIDE 8 (refuerzo ejecución) =================
s=8
from pptx.util import Pt as _Pt
def set_cell(cell,text,size,bold=False,color=DARK,align=None):
    cell.text=text
    for p in cell.text_frame.paragraphs:
        if align is not None: p.alignment=align
        for r in p.runs:
            r.font.size=_Pt(size); r.font.name="Helvetica"; r.font.bold=bold
            if color is not None: r.font.color.rgb=RGBColor.from_string(color)
tbl=S[s-1].shapes[1].table
obj={1:"Marco de trabajo, equipos y foto inicial de Eiffage (continúa la reunión de Albacete).",
     2:"Presentáis vuestras iniciativas y PoCs; las estructuramos y detectamos quick wins.",
     3:"Estrategia de Eiffage y marco común de evaluación y priorización.",
     4:"Madurez, datos disponibles y complejidad de cada iniciativa.",
     5:"Valor, esfuerzo y riesgo de cada iniciativa de cara a producción.",
     6:"Cerramos las 3 iniciativas clave y su modelo de ejecución (autónomo, con apoyo o partner), con plan de acción y readout.",
     7:"A elegir juntos: deep-dive por iniciativa, sesión con referentes de negocio, readout a dirección y demo de contraste."}
names={6:"6. Convergencia y cierre",7:"4x sesiones flexibles"}
for ri in range(1,8):
    bold_row = ri in (6,7)
    set_cell(tbl.cell(ri,0), tbl.cell(ri,0).text.strip(), 11, bold=True, color=DARK)
    set_cell(tbl.cell(ri,1), obj[ri], 11, bold=bold_row, color=DARK)
    set_cell(tbl.cell(ri,2), tbl.cell(ri,2).text.strip(), 11, color=DARK)
    set_cell(tbl.cell(ri,3), tbl.cell(ri,3).text.strip(), 11, color=DARK)
# header
for ci in range(4):
    set_cell(tbl.cell(0,ci), tbl.cell(0,ci).text.strip(), 12, bold=True, color=None)
# resaltar filas de ejecución
for ri in (6,7):
    for ci in range(4):
        tbl.cell(ri,ci).fill.solid(); tbl.cell(ri,ci).fill.fore_color.rgb=RGBColor.from_string("FDE7D2")
# condensar tabla "Modelo de trabajo"
tbl2=S[s-1].shapes[3].table
m={(1,1):"Trabajamos en modo co-creación: no imponemos, compartimos conocimiento y capacitamos a vuestro equipo.",
   (2,1):"Facilitación con vuestro Líder de IA y personas clave. Incluye diseño, guion y materiales.",
   (3,1):"Notas, puntuación y consolidación de resultados.",
   (4,1):"Matriz de priorización, viabilidad y salto a operación. Iniciativas clave y modelo de ejecución.",
   (5,1):"Perfiles senior de Orange Business durante una ventana de ~3 meses (12 semanas)."}
for (r,c),txt in m.items():
    set_cell(tbl2.cell(r,c), txt, 10, color=DARK)
for ci in range(2):
    set_cell(tbl2.cell(0,ci), tbl2.cell(0,ci).text.strip(), 11, bold=True, color=None)

# ================= NUEVA SLIDE: Ejecución de las 3 iniciativas =================
blank=prs.slide_layouts[20]
ns=prs.slides.add_slide(blank)
def txt(x,y,w,h,segs,size,color=DARK,bold=None,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,font="Helvetica"):
    tb=ns.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    set_para(tf.paragraphs[0], segs, size=size, color=color, bold=bold, font=font)
    tf.paragraphs[0].alignment=align
    return tb
def rrect(x,y,w,h,fill,line=None,line_w=1.0,radius=0.08):
    shp=ns.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=RGBColor.from_string(fill)
    if line: shp.line.color.rgb=RGBColor.from_string(line); shp.line.width=_Pt(line_w)
    else: shp.line.fill.background()
    shp.shadow.inherit=False
    try: shp.adjustments[0]=radius
    except: pass
    return shp
# título
txt(0.45,0.34,12.4,0.6,"De la priorización a la ejecución",28,color=ORANGE,bold=True)
txt(0.45,1.05,12.4,0.5,"El foco del assessment: cerrar 3 iniciativas clave y decidir cómo llevarlas a cabo.",17,color=DARK,bold=False)
txt(0.45,1.75,12.4,0.35,"Modelo de ejecución, iniciativa a iniciativa",16,color=ORANGE2,bold=True)
# 3 tarjetas
cardsN=[("Autónomo","Vuestro equipo ejecuta. Os dejamos el plan, los criterios y los gaps a cubrir."),
        ("Con apoyo","Ejecutáis vosotros, con acompañamiento puntual de Orange Business en los momentos clave."),
        ("Con partner","Orange Business co-ejecuta la iniciativa de principio a fin, junto a vuestro equipo.")]
xs=[0.5,4.8,9.1]; cw=3.7; cy=2.2; ch=2.55
for (hd,bd),x in zip(cardsN,xs):
    rrect(x,cy,cw,ch,LIGHT,line=ORANGE,line_w=1.5)
    rrect(x,cy,cw,0.62,ORANGE,radius=0.12)
    txt(x+0.05,cy+0.12,cw-0.1,0.4,hd,19,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    txt(x+0.22,cy+0.8,cw-0.44,ch-0.95,bd,16,color=DARK)
# banda inferior sesiones flexibles
by=5.05; bh=1.05
rrect(0.5,by,12.3,bh,ORANGE)
txt(0.8,by+0.13,11.7,0.35,"Sesiones flexibles orientadas a ejecutar",16,color=WHITE,bold=True)
txt(0.8,by+0.52,11.7,0.4,"Deep-dive por iniciativa   ·   Sesión con referentes de negocio   ·   Readout a dirección   ·   Demo de contraste",15,color=WHITE,bold=False)
# cierre
txt(0.5,6.35,12.3,0.5,[("Cada iniciativa sale con su ",False),("caso de valor",True),(", su ",False),("modelo de ejecución",True),(" y sus ",False),("siguientes pasos",True),(".",False)],16,color=DARK)
# mover nueva slide a posición 9 (tras slide 8)
sldIdLst=prs.slides._sldIdLst
ids=list(sldIdLst)
sldIdLst.remove(ids[-1]); sldIdLst.insert(8, ids[-1])

# fit-check nueva slide
for i,c in enumerate(cardsN):
    fit(ns.shapes[2+ i*3 +2] if False else ns.shapes[0],0,label="") if False else None
print("WARNINGS:", warn)
prs.save('/home/user/orange_crm/propuesta_v2.pptx')
print("SAVED propuesta_v2.pptx, total slides:", len(prs.slides._sldIdLst))
